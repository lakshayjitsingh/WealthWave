from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy
from lxml import etree

# === LOAD THE TEMPLATE TO COPY ITS BACKGROUND ===
TEMPLATE_PATH = r"C:\Users\Owner\Downloads\Project PPT_Format.pptx"
OUTPUT_PATH = r"C:\Users\Owner\Downloads\WealthWave_Presentation.pptx"

template = Presentation(TEMPLATE_PATH)
prs = Presentation(TEMPLATE_PATH)  # We build ON TOP of the template

# Get slide dimensions
W = prs.slide_width
H = prs.slide_height

# === COLOR PALETTE (WealthWave Brand) ===
PURPLE      = RGBColor(0x7C, 0x3A, 0xED)   # accent purple
LIGHT_PURPLE= RGBColor(0xA7, 0x8B, 0xFA)   # lighter purple
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xE2, 0xE8, 0xF0)
DARK_TEXT   = RGBColor(0x1E, 0x29, 0x3B)
ACCENT      = RGBColor(0x06, 0xB6, 0xD4)   # cyan accent

def copy_background(src_slide, dst_slide):
    """Copy the background XML from source slide to destination slide."""
    src_bg = src_slide.background
    dst_bg = dst_slide.background
    src_fill = src_bg.fill
    dst_fill = dst_bg.fill
    # Copy the background element
    src_bg_elem = src_slide._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}bg')
    spTree_src = src_slide._element.spTree
    spTree_dst = dst_slide._element.spTree

def get_template_slide_bg(template, slide_index=0):
    """Returns the background XML element of the template slide."""
    return template.slides[slide_index]._element

def add_slide_with_bg(prs, template):
    """Add a blank slide and copy the background from the first template slide."""
    blank_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(blank_layout)
    
    # Copy background from template slide 0
    template_slide_elem = template.slides[0]._element
    new_slide_elem = slide._element
    
    # Copy the cSld (common slide data) background
    t_cSld = template_slide_elem.find('{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
    n_cSld = new_slide_elem.find('{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
    
    if t_cSld is not None and n_cSld is not None:
        t_bg = t_cSld.find('{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
        if t_bg is not None:
            t_bg_copy = copy.deepcopy(t_bg)
            # Remove existing bg if any
            n_bg = n_cSld.find('{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
            if n_bg is not None:
                n_cSld.remove(n_bg)
            n_cSld.insert(0, t_bg_copy)
    
    return slide

def add_textbox(slide, text, left, top, width, height, font_size=18, bold=False, 
                color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_colored_box(slide, left, top, width, height, color, transparency=0):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_bullet_textbox(slide, bullets, left, top, width, height, font_size=14, color=WHITE, icon=""):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"{icon}  {bullet}" if icon else f"▸  {bullet}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox

# ============================================================
# Remove all existing slides from the copied presentation
# ============================================================
prs2 = Presentation(TEMPLATE_PATH)

# Delete all existing slides cleanly
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
NSMAP = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
slide_ids = list(prs2.slides._sldIdLst)
for sld_id in slide_ids:
    rId = sld_id.get('{%s}id' % NSMAP)
    if rId:
        prs2.part.drop_rel(rId)
    prs2.slides._sldIdLst.remove(sld_id)

# ============================================================
# SLIDE 1: TITLE SLIDE
# ============================================================
slide1 = add_slide_with_bg(prs2, template)

# Purple accent bar on left
add_colored_box(slide1, Inches(0), Inches(0), Inches(0.15), H, PURPLE)

# WealthWave Logo text
add_textbox(slide1, "WealthWave AI", Inches(0.4), Inches(1.2), Inches(8), Inches(1.5),
            font_size=48, bold=True, color=PURPLE, align=PP_ALIGN.LEFT)

# Subtitle
add_textbox(slide1, "A Cloud-Native Financial Co-pilot with Real-Time AI Insights",
            Inches(0.4), Inches(2.7), Inches(8), Inches(1),
            font_size=20, bold=False, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

# Divider line shape
line = slide1.shapes.add_shape(1, Inches(0.4), Inches(3.8), Inches(4), Inches(0.05))
line.fill.solid()
line.fill.fore_color.rgb = PURPLE
line.line.fill.background()

# Student info
add_textbox(slide1, "Presented by: Lakshayjit Singh  |  24BAI70225",
            Inches(0.4), Inches(4.1), Inches(8), Inches(0.5),
            font_size=14, color=LIGHT_GRAY)
add_textbox(slide1, "Capstone Project  |  B.Tech AI & ML  |  2024–25",
            Inches(0.4), Inches(4.6), Inches(8), Inches(0.5),
            font_size=12, color=RGBColor(0xA0, 0xA0, 0xA0))

# ============================================================
# SLIDE 2: THE PROBLEM
# ============================================================
slide2 = add_slide_with_bg(prs2, template)
add_colored_box(slide2, Inches(0), Inches(0), Inches(0.15), H, PURPLE)
add_textbox(slide2, "The Problem", Inches(0.4), Inches(0.3), Inches(9), Inches(0.8),
            font_size=36, bold=True, color=PURPLE)
add_colored_box(slide2, Inches(0.4), Inches(1.15), Inches(4), Inches(0.05), PURPLE)

problems = [
    "Financial data stored locally (SQLite) is inaccessible across devices",
    "Users lack real-time AI-driven analysis of their spending habits",
    "Desktop apps and web apps often have 'Data Gaps' — information doesn't sync",
    "No unified dashboard for income tracking, tax estimation, and savings goals",
    "Traditional finance apps are complex, ugly, and not beginner-friendly",
]
add_bullet_textbox(slide2, problems, Inches(0.5), Inches(1.4), Inches(8.5), Inches(4.5),
                   font_size=16, icon="❌")

# ============================================================
# SLIDE 3: OUR SOLUTION
# ============================================================
slide3 = add_slide_with_bg(prs2, template)
add_colored_box(slide3, Inches(0), Inches(0), Inches(0.15), H, PURPLE)
add_textbox(slide3, "Our Solution", Inches(0.4), Inches(0.3), Inches(9), Inches(0.8),
            font_size=36, bold=True, color=PURPLE)
add_colored_box(slide3, Inches(0.4), Inches(1.15), Inches(4), Inches(0.05), PURPLE)

solutions = [
    "Cloud-First Architecture: Unified data on Neon PostgreSQL — accessible everywhere",
    "AI Integration: Google Gemini-2.0 via OpenRouter for smart financial coaching",
    "Dual-Platform: Seamless experience on Web (Vercel) and Windows Desktop (Electron)",
    "Beautiful Glassmorphism UI with Dark/Light mode and real-time animated charts",
    "One-Click PDF & CSV Report Export, Multi-Currency & Tax Estimator built-in",
]
add_bullet_textbox(slide3, solutions, Inches(0.5), Inches(1.4), Inches(8.5), Inches(4.5),
                   font_size=16, icon="✅")

# ============================================================
# SLIDE 4: TECHNOLOGY STACK
# ============================================================
slide4 = add_slide_with_bg(prs2, template)
add_colored_box(slide4, Inches(0), Inches(0), Inches(0.15), H, PURPLE)
add_textbox(slide4, "Technology Stack", Inches(0.4), Inches(0.3), Inches(9), Inches(0.8),
            font_size=36, bold=True, color=PURPLE)
add_colored_box(slide4, Inches(0.4), Inches(1.15), Inches(4), Inches(0.05), PURPLE)

# 2 Column Layout
col1 = [("⚛️  Frontend", "React 19, Vite, React Router, Recharts"),
        ("🎨  Styling", "Vanilla CSS, Glassmorphism, Animations"),
        ("🖥️  Desktop App", "Electron.js (Windows .exe Installer)")]
col2 = [("🟢  Backend", "Node.js, Express.js, JWT Auth, Bcrypt"),
        ("🐘  Database", "PostgreSQL on Neon Cloud"),
        ("☁️  Deployment", "Vercel (Frontend) + Render (Backend)"),
        ("🤖  AI Engine", "Google Gemini-2.0 via OpenRouter API")]

for i, (title, val) in enumerate(col1):
    y = Inches(1.5 + i * 1.1)
    add_colored_box(slide4, Inches(0.4), y, Inches(4.1), Inches(0.9), RGBColor(0x2D, 0x1B, 0x69))
    add_textbox(slide4, title, Inches(0.55), y + Pt(4), Inches(3.8), Inches(0.4),
                font_size=13, bold=True, color=LIGHT_PURPLE)
    add_textbox(slide4, val, Inches(0.55), y + Pt(22), Inches(3.8), Inches(0.4),
                font_size=12, color=WHITE)

for i, (title, val) in enumerate(col2):
    y = Inches(1.5 + i * 0.85)
    add_colored_box(slide4, Inches(4.8), y, Inches(4.6), Inches(0.75), RGBColor(0x0C, 0x3A, 0x4A))
    add_textbox(slide4, title, Inches(4.95), y + Pt(4), Inches(4.2), Inches(0.35),
                font_size=13, bold=True, color=ACCENT)
    add_textbox(slide4, val, Inches(4.95), y + Pt(22), Inches(4.2), Inches(0.35),
                font_size=11, color=WHITE)

# ============================================================
# SLIDE 5: CLOUD ARCHITECTURE
# ============================================================
slide5 = add_slide_with_bg(prs2, template)
add_colored_box(slide5, Inches(0), Inches(0), Inches(0.15), H, PURPLE)
add_textbox(slide5, "Cloud Architecture", Inches(0.4), Inches(0.3), Inches(9), Inches(0.8),
            font_size=36, bold=True, color=PURPLE)
add_colored_box(slide5, Inches(0.4), Inches(1.15), Inches(4), Inches(0.05), PURPLE)

arch_items = [
    ("User (Browser / Windows App)", LIGHT_PURPLE, Inches(3.5), Inches(1.3)),
    ("▼  HTTPS Request", WHITE, Inches(4.2), Inches(2.0)),
    ("Vercel  →  React Frontend", RGBColor(0x00, 0xC7, 0xBE), Inches(3.2), Inches(2.5)),
    ("▼  REST API Call", WHITE, Inches(4.2), Inches(3.2)),
    ("Render  →  Node.js Backend + AI", RGBColor(0x46, 0xC8, 0x7F), Inches(3.0), Inches(3.7)),
    ("▼  SQL Query", WHITE, Inches(4.2), Inches(4.4)),
    ("Neon  →  PostgreSQL Cloud DB", RGBColor(0x00, 0x78, 0xD4), Inches(3.0), Inches(4.9)),
]

for text, color, left, top in arch_items:
    add_textbox(slide5, text, left, top, Inches(4), Inches(0.5),
                font_size=15 if "▼" not in text else 12,
                bold="▼" not in text,
                color=color, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 6: KEY FEATURES
# ============================================================
slide6 = add_slide_with_bg(prs2, template)
add_colored_box(slide6, Inches(0), Inches(0), Inches(0.15), H, PURPLE)
add_textbox(slide6, "Key Features", Inches(0.4), Inches(0.3), Inches(9), Inches(0.8),
            font_size=36, bold=True, color=PURPLE)
add_colored_box(slide6, Inches(0.4), Inches(1.15), Inches(4), Inches(0.05), PURPLE)

features = [
    ("🤖 WealthWave AI", "Google Gemini-2.0 powered chatbot with full financial context awareness"),
    ("📊 Smart Dashboard", "Real-time charts: Cash Flow bar charts + Expense Breakdown pie charts"),
    ("🌍 Multi-Currency", "Live exchange rates for INR, USD, GBP, EUR with automatic conversion"),
    ("🧾 Tax Estimator", "Dynamic tax calculation for India, USA, UK & Germany"),
    ("📥 Report Export", "One-click PDF & CSV export of all transactions and analytics"),
    ("🏆 Savings Goals", "Track progress towards specific financial goals with visual meters"),
]

for i, (title, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    x = Inches(0.4 + col * 4.8)
    y = Inches(1.45 + row * 1.25)
    add_colored_box(slide6, x, y, Inches(4.4), Inches(1.1), RGBColor(0x1E, 0x10, 0x4A))
    add_textbox(slide6, title, x + Inches(0.1), y + Pt(5), Inches(4.1), Inches(0.35),
                font_size=13, bold=True, color=LIGHT_PURPLE)
    add_textbox(slide6, desc, x + Inches(0.1), y + Pt(25), Inches(4.1), Inches(0.5),
                font_size=11, color=LIGHT_GRAY)

# ============================================================
# SLIDE 7: LIVE DEMO SCREENSHOTS (Placeholder text)
# ============================================================
slide7 = add_slide_with_bg(prs2, template)
add_colored_box(slide7, Inches(0), Inches(0), Inches(0.15), H, PURPLE)
add_textbox(slide7, "Live Application", Inches(0.4), Inches(0.3), Inches(9), Inches(0.8),
            font_size=36, bold=True, color=PURPLE)
add_colored_box(slide7, Inches(0.4), Inches(1.15), Inches(4), Inches(0.05), PURPLE)

# Placeholder boxes for screenshots
box1 = add_colored_box(slide7, Inches(0.4), Inches(1.4), Inches(4.3), Inches(3.2), RGBColor(0x1A, 0x1A, 0x2E))
add_textbox(slide7, "[ Login Page Screenshot ]\nwealth-wave-gamma.vercel.app",
            Inches(0.5), Inches(2.5), Inches(4), Inches(1),
            font_size=13, color=LIGHT_PURPLE, align=PP_ALIGN.CENTER)

box2 = add_colored_box(slide7, Inches(4.9), Inches(1.4), Inches(4.3), Inches(3.2), RGBColor(0x1A, 0x1A, 0x2E))
add_textbox(slide7, "[ Dashboard Screenshot ]\nReal-time Charts & AI Panel",
            Inches(5.0), Inches(2.5), Inches(4), Inches(1),
            font_size=13, color=ACCENT, align=PP_ALIGN.CENTER)

add_textbox(slide7, "🌐  Live URL:  wealth-wave-gamma.vercel.app",
            Inches(0.4), Inches(4.8), Inches(9), Inches(0.4),
            font_size=14, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 8: SECURITY & DATABASE
# ============================================================
slide8 = add_slide_with_bg(prs2, template)
add_colored_box(slide8, Inches(0), Inches(0), Inches(0.15), H, PURPLE)
add_textbox(slide8, "Security & Database", Inches(0.4), Inches(0.3), Inches(9), Inches(0.8),
            font_size=36, bold=True, color=PURPLE)
add_colored_box(slide8, Inches(0.4), Inches(1.15), Inches(4), Inches(0.05), PURPLE)

sec = [
    "🔐  JWT Authentication: Stateless, token-based auth for every API request",
    "🔒  Bcrypt Password Hashing: Passwords are NEVER stored in plain text",
    "🛡️  SQL Parameterization: All queries use $1, $2 params to prevent SQL Injection",
    "🌐  CORS Policy: Backend only accepts requests from the trusted Vercel domain",
    "☁️  Neon PostgreSQL: Cloud database with SSL encryption on all connections",
    "🔑  Environment Variables: API keys and secrets stored securely on Render — never in code",
]
add_bullet_textbox(slide8, sec, Inches(0.5), Inches(1.4), Inches(8.8), Inches(5),
                   font_size=15, icon="")

# ============================================================
# SLIDE 9: FUTURE SCOPE
# ============================================================
slide9 = add_slide_with_bg(prs2, template)
add_colored_box(slide9, Inches(0), Inches(0), Inches(0.15), H, PURPLE)
add_textbox(slide9, "Future Scope", Inches(0.4), Inches(0.3), Inches(9), Inches(0.8),
            font_size=36, bold=True, color=PURPLE)
add_colored_box(slide9, Inches(0.4), Inches(1.15), Inches(4), Inches(0.05), PURPLE)

future = [
    "📱  Android & iOS Mobile App using React Native with the same cloud backend",
    "🏦  UPI / Bank Integration: Auto-import transactions from bank statements",
    "📈  Investment Portfolio Tracker: Stocks, Mutual Funds & Crypto in one place",
    "🧠  Predictive AI: Machine learning model to forecast future expenses and savings",
    "👨‍👩‍👧  Family Budgeting Mode: Shared financial dashboard for household management",
    "🔔  Smart Alerts: Push notifications for unusual spending and budget thresholds",
]
add_bullet_textbox(slide9, future, Inches(0.5), Inches(1.4), Inches(8.8), Inches(5),
                   font_size=15, icon="")

# ============================================================
# SLIDE 10: CONCLUSION / THANK YOU
# ============================================================
slide10 = add_slide_with_bg(prs2, template)
add_colored_box(slide10, Inches(0), Inches(0), Inches(0.15), H, PURPLE)

add_textbox(slide10, "Thank You!", Inches(0.4), Inches(1.0), Inches(9), Inches(1.2),
            font_size=52, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

add_textbox(slide10, "WealthWave AI — Mastering Money, Anywhere, Anytime.",
            Inches(0.4), Inches(2.4), Inches(9), Inches(0.7),
            font_size=18, italic=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_colored_box(slide10, Inches(3), Inches(3.2), Inches(3.5), Inches(0.05), PURPLE)

info = [
    "🌐  wealth-wave-gamma.vercel.app",
    "👨‍💻  GitHub: github.com/lakshayjitsingh/WealthWave",
    "🎓  Lakshayjit Singh  |  24BAI70225",
]
for i, line in enumerate(info):
    add_textbox(slide10, line, Inches(0.4), Inches(3.5 + i*0.55), Inches(9), Inches(0.5),
                font_size=15, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# SAVE
# ============================================================
prs2.save(OUTPUT_PATH)
print("SUCCESS! PPT saved to: " + OUTPUT_PATH)
